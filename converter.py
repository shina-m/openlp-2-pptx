import json
import os
import sys
import random
import zipfile
import shutil

from pathlib import Path

from xml.dom import minidom
from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_BLANK_PAGE = 6
R16_9_INCHES = (10, 5.625)
R16_10_INCHES = (10, 6.25)
PRS_SIZE = tuple(map(Inches, R16_10_INCHES))

SLIDE_MARGIN = Inches(0.2)
CONTENT_LEFT = SLIDE_MARGIN
CONTENT_TOP = SLIDE_MARGIN
CONTENT_WIDTH = PRS_SIZE[0] - (2 * SLIDE_MARGIN)
CONTENT_HEIGHT = PRS_SIZE[1] - (4 * SLIDE_MARGIN)
FOOTER_TOP = CONTENT_TOP + CONTENT_HEIGHT
FOOTER_HEIGHT = 2 * SLIDE_MARGIN
P_LEFT = P_TOP = Inches(0)

DEFAULT_THEME = "light"
BKG_IMAGES = {
    "dark": ["dark_" + str(x) for x in range(4)],
    "light": ["light_" + str(x) for x in range(4)],
}

if getattr(sys, 'frozen', False):
        bundle_dir = sys._MEIPASS
else:
        bundle_dir = os.path.dirname(os.path.abspath(__file__))

class Slide:
    def __init__(self):
        self.id = None
        self.name = None
        self.lines = []

class Song:
    def __init__(self):
        self.title = None
        self.author = None
        self.slides = []
        self.verse_order = None
        self.bkg_image_path = None

class Job:
    def __init__(self, osz_file_path, theme):
        self.osz_file_path = osz_file_path
        self.theme = theme
        self.prs_list = []
        self.songs = []
        self.default_output_path = None
        images = BKG_IMAGES.get(theme, DEFAULT_THEME)
        self.bkg_images = random.sample(images, len(images))


    def _xml_to_song(self, song_xml):
        verses_chorus_map = {
            'v': "Verse",
            'c': "Chorus",
            'b': "Bridge"
        }

        song = Song()
        song.title = song_xml.getElementsByTagName('title')[0].firstChild.data
        song.author = song_xml.getElementsByTagName('author')[0].firstChild.data

        verseOrder = song_xml.getElementsByTagName('verseOrder')
        if len(verseOrder) != 0:
            song.verse_order = verseOrder[0].firstChild.data.split(" ")

        for child in song_xml.getElementsByTagName('lyrics')[0].childNodes:
            # if child.localName != 'verse':
            #     continue

            # get slide name
            vc = child.attributes['name'].value
            # create new slide
            slide = Slide()
            slide.id = vc
            song.slides.append(slide)

            slide.name = verses_chorus_map.get(vc[0], vc[0]) + ' ' + vc[1]
            if len(vc) == 3:
                slide.name += vc[2]

            # get content for the slide
            for line in child.getElementsByTagName('lines'):
                words = line
                tags = child.getElementsByTagName('tag')
                if len(tags) == 1:
                    words.childNodes = tags[0].childNodes + words.childNodes

                for l in words.childNodes:
                    if l.nodeValue:
                        slide.lines.append(l.nodeValue)
        return song

    def get_songs_from_osz(self):
        items = None
        try:
            with zipfile.ZipFile(self.osz_file_path, "r") as zip_ref:
                zip_ref.extractall(bundle_dir)
                osj_file_path = "{}/{}".format(bundle_dir, zip_ref.namelist()[0])
                with open(osj_file_path, 'r') as file_to:
                    if osj_file_path.endswith('osj'):
                        items = json.load(file_to)
                    else:
                        print('The service file you are trying to open is in an old')
            os.remove(osj_file_path)
        except Exception as e:
            if hasattr(e, 'message'):
                return e.message
            else:
                return e

        for item in items:
            if 'serviceitem' in item \
                    and 'header' in item['serviceitem'] \
                    and item['serviceitem']['header']['name'] == "songs":
                xml_string = item['serviceitem']['header'].get('xml_version', None)
                xml = minidom.parseString(xml_string)
                if xml:
                    song_obj = self._xml_to_song(xml)
                    image_file_name = self.bkg_images[len(self.songs) % len(self.bkg_images)] + ".jpg"
                    song_obj.bkg_image_path = os.path.join(bundle_dir, "images/", image_file_name)
                    self.songs.append(song_obj)

    def _add_title_slide(self, prs, song_index):
        song = self.songs[song_index]
        slide = prs.slides.add_slide(slide_layout=prs.slide_layouts[SLD_LAYOUT_TITLE])
        #     slide.shapes.title.text = song.title
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.text = song.title

        rgb = RGBColor(128, 0, 0) if self.theme == "light" else RGBColor(225, 225, 225)
        p.font.color.rgb = rgb
        p.font.name = 'Gill Sans'  # 'Myriad Pro'
        p.font.bold = True
        p.font.size = Pt(40)

        pic = slide.shapes.add_picture(song.bkg_image_path, P_LEFT, P_TOP, width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    def _add_footer(self, song: Song, prs_slide, slide_loc_str):
        color = MSO_THEME_COLOR_INDEX.DARK_1 if self.theme == "light" else MSO_THEME_COLOR_INDEX.LIGHT_1

        txBox = prs_slide.shapes.add_textbox(left=CONTENT_LEFT, top=FOOTER_TOP, width=CONTENT_WIDTH / 2, height=FOOTER_HEIGHT)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.font.italic = True
        p.font.size = Pt(12)
        p.font.color.theme_color = color
        p.text = prs_slide.name

        p = tf.add_paragraph()
        p.font.italic = True
        p.font.size = Pt(11)
        p.font.color.theme_color = color
        slide_loc_str_split = slide_loc_str.split("/")
        if slide_loc_str_split[0] == slide_loc_str_split[1]:
            p.font.bold = True
            p.font.color.rgb = RGBColor(128, 0, 0)

        p.text = "({})".format(slide_loc_str)

        txBox = prs_slide.shapes.add_textbox(left=(CONTENT_WIDTH / 2) + (SLIDE_MARGIN), top=FOOTER_TOP, width=CONTENT_WIDTH / 2,
                                             height=FOOTER_HEIGHT)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.font.italic = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        p.font.size = Pt(12)
        p.font.color.theme_color = color
        p.text = "Title: {} \nAuthor: {}".format(song.title, song.author)

    def _add_prs_slide(self, prs, song, slide, slide_loc_str):
        prs_slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK_PAGE])

        # Add bkg picture
        pic = prs_slide.shapes.add_picture(song.bkg_image_path, P_LEFT, P_TOP, width=prs.slide_width, height=prs.slide_height)
        prs_slide.shapes._spTree.remove(pic._element)
        prs_slide.shapes._spTree.insert(2, pic._element)

        prs_slide.name = slide.name
        self._add_footer(song, prs_slide, slide_loc_str)

        # create text box
        txBox = prs_slide.shapes.add_textbox(left=CONTENT_LEFT, top=CONTENT_TOP, width=CONTENT_WIDTH, height=CONTENT_HEIGHT)
        tf = txBox.text_frame
        for i, line in enumerate(slide.lines):
            p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.text = line
            p.font.name = 'Gill Sans'  # 'Myriad Pro'
            p.font.size = Pt(40)
            color = MSO_THEME_COLOR_INDEX.DARK_1 if self.theme == "light" else MSO_THEME_COLOR_INDEX.LIGHT_1
            p.font.color.theme_color = color
            if "Chorus" in slide.name:
                p.font.italic = True

        # format textbox
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        # tf.fit_text(font_family='Gill Sans', max_size=40, font_file="/System/Library/Fonts/Supplemental/GillSans.ttc")

    def generate_ppt(self):
        for i, song in enumerate(self.songs):
            prs = Presentation()
            prs.slide_width, prs.slide_height = PRS_SIZE
            self._add_title_slide(prs, i)

            slides = []
            if song.verse_order is None:
                for s in song.slides:
                    slides.append(s)
            else:
                for v in song.verse_order:
                    for s in song.slides:
                        if v == s.id[:2]:
                            slides.append(s)

            for i, s in enumerate(slides):
                slide_loc_str = "{}/{}".format(i+1, len(slides))
                self._add_prs_slide(prs, song, s, slide_loc_str)

            self.prs_list.append({
                "title": song.title,
                "prs": prs,
            })
        # self.default_output_path = Path(self.osz_file_path).stem + ".pptx"

    def save_file(self, output_dir="output"):
        def cleanFilename(sourcestring, removestring="%:/,.\\[]<>*?"):
            return ''.join([c for c in sourcestring if c not in removestring])


        if os.path.exists(output_dir):
            shutil.rmtree(output_dir)
        os.makedirs(output_dir)
        for out in self.prs_list:
            filename = cleanFilename(out["title"])
            try:
                out["prs"].save("{}/{}.pptx".format(output_dir,filename))
            except Exception as e:
                if hasattr(e, 'message'):
                    print(e.message)
                else:
                    print(e)



def gui_endpoint(osz_file_path, ppt_file_path=""):
    job = Job(osz_file_path=osz_file_path, theme="light")
    err = job.get_songs_from_osz()
    if err:
        return err
    job.generate_ppt()
    err = job.save_file(ppt_file_path)
    if err:
        return err
